"""
文本抽取工具模块
支持从多种文档格式中抽取纯文本内容
"""
import json
import logging
import mimetypes
from pathlib import Path
from typing import Optional, Tuple

logger = logging.getLogger(__name__)


class TextExtractionError(Exception):
    """文本抽取异常"""
    pass


def extract_text_from_file(file_path: str, mime_type: Optional[str] = None) -> str:
    """
    从文件中抽取文本内容
    
    Args:
        file_path: 文件路径
        mime_type: MIME类型（可选，用于辅助判断）
    
    Returns:
        抽取的文本内容 (UTF-8)
    
    Raises:
        TextExtractionError: 当抽取失败时
    """
    path = Path(file_path)
    
    if not path.exists():
        raise TextExtractionError(f"文件不存在: {file_path}")
    
    # 获取文件扩展名
    ext = path.suffix.lower()
    
    # 如果没有提供 mime_type，尝试从文件名推断
    if not mime_type:
        mime_type, _ = mimetypes.guess_type(file_path)
    
    try:
        # 根据扩展名选择合适的抽取方法
        if ext in ['.txt', '.md']:
            return _extract_text(file_path)
        elif ext == '.json':
            return _extract_json(file_path)
        elif ext == '.csv':
            return _extract_csv(file_path)
        elif ext == '.pdf':
            return _extract_pdf(file_path)
        elif ext in ['.doc', '.docx']:
            return _extract_docx(file_path)
        elif ext in ['.ppt', '.pptx']:
            return _extract_pptx(file_path)
        else:
            # 未知类型，尝试作为纯文本读取
            logger.warning(f"未知文件类型 {ext}，尝试作为纯文本读取")
            return _extract_text(file_path)
    except Exception as e:
        logger.error(f"文本抽取失败 file={file_path} error={str(e)}", exc_info=True)
        raise TextExtractionError(f"无法从该文件抽取文本（{ext}）: {str(e)}")


def _extract_text(file_path: str) -> str:
    """抽取纯文本文件"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # 尝试其他编码
        try:
            with open(file_path, 'r', encoding='gbk') as f:
                return f.read()
        except Exception:
            with open(file_path, 'r', encoding='latin-1') as f:
                return f.read()


def _extract_json(file_path: str) -> str:
    """抽取JSON文件并格式化"""
    with open(file_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    # 格式化输出，便于阅读
    return json.dumps(data, ensure_ascii=False, indent=2)


def _extract_csv(file_path: str) -> str:
    """抽取CSV文件（保持原始格式）"""
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()


def _extract_pdf(file_path: str) -> str:
    """从PDF抽取文本"""
    try:
        from pypdf import PdfReader
    except ImportError:
        raise TextExtractionError("缺少 pypdf 库，无法处理 PDF 文件")
    
    try:
        reader = PdfReader(file_path)
        text_parts = []
        
        for i, page in enumerate(reader.pages):
            try:
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(f"=== 第 {i + 1} 页 ===\n{page_text}")
            except Exception as e:
                logger.warning(f"PDF页面 {i+1} 抽取失败: {e}")
                continue
        
        if not text_parts:
            raise TextExtractionError("PDF 文件中未能抽取到任何文本（可能是扫描件或图片PDF）")
        
        return "\n\n".join(text_parts)
    except Exception as e:
        if isinstance(e, TextExtractionError):
            raise
        raise TextExtractionError(f"PDF 处理失败: {str(e)}")


def _extract_docx(file_path: str) -> str:
    """从DOCX抽取文本"""
    try:
        from docx import Document
    except ImportError:
        raise TextExtractionError("缺少 python-docx 库，无法处理 DOCX 文件")
    
    try:
        doc = Document(file_path)
        text_parts = []
        
        # 抽取段落
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # 抽取表格
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_text.append(row_text)
            if table_text:
                text_parts.append("\n[表格]\n" + "\n".join(table_text))
        
        if not text_parts:
            raise TextExtractionError("DOCX 文件中未能抽取到任何文本")
        
        return "\n\n".join(text_parts)
    except Exception as e:
        if isinstance(e, TextExtractionError):
            raise
        raise TextExtractionError(f"DOCX 处理失败: {str(e)}")


def _extract_pptx(file_path: str) -> str:
    """从PPTX抽取文本"""
    try:
        from pptx import Presentation
    except ImportError:
        raise TextExtractionError("缺少 python-pptx 库，无法处理 PPTX 文件")
    
    try:
        prs = Presentation(file_path)
        text_parts = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_text.append(f"=== 幻灯片 {i + 1} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:  # 除了标题外还有内容
                text_parts.append("\n".join(slide_text))
        
        if not text_parts:
            raise TextExtractionError("PPTX 文件中未能抽取到任何文本")
        
        return "\n\n".join(text_parts)
    except Exception as e:
        if isinstance(e, TextExtractionError):
            raise
        raise TextExtractionError(f"PPTX 处理失败: {str(e)}")


def get_safe_filename(original_name: str) -> str:
    """
    生成安全的文件名，防止路径穿越
    
    Args:
        original_name: 原始文件名
    
    Returns:
        安全的文件名
    """
    # 移除路径分隔符
    safe_name = original_name.replace('/', '_').replace('\\', '_')
    # 移除其他潜在危险字符
    safe_name = safe_name.replace('..', '_')
    # 限制长度
    if len(safe_name) > 255:
        # 保留扩展名
        name_part = safe_name[:200]
        ext_part = Path(safe_name).suffix[-50:] if Path(safe_name).suffix else ""
        safe_name = name_part + ext_part
    return safe_name


def get_file_extension(filename: str) -> str:
    """获取文件扩展名（小写，包含点）"""
    return Path(filename).suffix.lower()


def is_allowed_extension(filename: str, allowed_extensions: set[str]) -> bool:
    """
    检查文件扩展名是否允许
    
    Args:
        filename: 文件名
        allowed_extensions: 允许的扩展名集合（如 {'.txt', '.pdf'}）
    
    Returns:
        是否允许
    """
    ext = get_file_extension(filename)
    return ext in allowed_extensions
